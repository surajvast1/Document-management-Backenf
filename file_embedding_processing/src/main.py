import logging
import json
import boto3
import pandas as pd
import os
import io
from qdrant_client import QdrantClient, models
from langchain_qdrant import QdrantVectorStore
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.embeddings import OpenAIEmbeddings
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

# Setting up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger()

# AWS S3 client setup
s3 = boto3.client('s3')

# Qdrant and OpenAI setup
QDRANT_URL = os.getenv('QDRANT_URL')
openai_key = os.getenv('OPENAI_API_KEY')
openai_org = os.getenv('OPENAI_ORGANIZATION')

embedding_model = OpenAIEmbeddings(openai_api_key=openai_key, openai_organization=openai_org)
client = QdrantClient(url=QDRANT_URL)

def collection_exists_safe(client, collection_name):
    try:
        collections_response = client.get_collections()
        collection_names = [collection.name for collection in collections_response.collections]
        return collection_name in collection_names
    except Exception as e:
        logger.error(f"Error checking collection existence: {e}")
        return False

def create_vectorstore(collection_name):
    if not collection_exists_safe(client, collection_name):
        try:
            client.create_collection(
                collection_name=collection_name,
                vectors_config=models.VectorParams(size=1536, distance=models.Distance.COSINE),
            )
            logger.info(f"Collection '{collection_name}' created successfully.")
        except Exception as e:
            logger.error(f"Failed to create collection: {e}")
            raise

    vectorstore = QdrantVectorStore(client=client, collection_name=collection_name, embedding=embedding_model)
    return vectorstore

def chunk_text(text, chunk_size=1000, overlap=100):
    splitter = CharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=overlap)
    return splitter.split_text(text)

def generate_embeddings(text):
    chunks = chunk_text(text)
    embeddings = [embedding_model.embed_documents([chunk])[0] for chunk in chunks]
    return embeddings, chunks

def read_file_from_s3(bucket, s3_key):
    try:
        obj = s3.get_object(Bucket=bucket, Key=s3_key)
        file_stream = io.BytesIO(obj['Body'].read())
        
        if s3_key.endswith('.csv'):
            df = pd.read_csv(file_stream)
            text = df.to_string(index=False)
        elif s3_key.endswith('.xlsx') or s3_key.endswith('.xls'):
            df = pd.read_excel(file_stream)
            text = df.to_string(index=False)
        elif s3_key.endswith('.pdf'):
            pdf_reader = PdfReader(file_stream)
            text = " ".join(page.extract_text() for page in pdf_reader.pages)
        elif s3_key.endswith('.pptx'):
            presentation = Presentation(file_stream)
            text = " ".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif s3_key.endswith('.docx'):
            doc = Document(file_stream)
            text = " ".join(para.text for para in doc.paragraphs)
        else:
            logger.warning(f"Unsupported file type: {s3_key}")
            return None
        
        logger.info(f"Successfully read file {s3_key}")
        return text
    except Exception as e:
        logger.error(f"Error reading file {s3_key}: {e}")
        return None

def store_embeddings_in_qdrant(vectorstore, embeddings, chunks, metadata):
    logger.info("Storing combined embeddings in Qdrant")
    payloads = []
    
    for i, (embedding, chunk) in enumerate(zip(embeddings, chunks)):
        metadata_with_chunk = metadata.copy()
        metadata_with_chunk['chunk_index'] = i
        metadata_with_chunk['text'] = chunk
        payloads.append(metadata_with_chunk)
    
    try:
        vectorstore.client.upload_collection(
            collection_name=vectorstore.collection_name,
            vectors=embeddings,
            payload=payloads
        )
        logger.info("Embeddings stored successfully")
    except Exception as e:
        logger.error(f"Error storing embeddings in Qdrant: {e}")
        raise

def lambda_handler(event, context):
    try:
        body = json.loads(event['body']) if 'body' in event and event['body'] else event
        
        bucket_name = body['bucket_name']
        user_id = body['user_id']
        context_id = body['context_id']
        name = body['name']
        
        # Construct S3 prefix and list files
        collection_prefix = f"{user_id}/{context_id}/{name}"
        response = s3.list_objects_v2(Bucket=bucket_name, Prefix=collection_prefix)
        
        if 'Contents' not in response:
            return {
                "statusCode": 404,
                "body": json.dumps({"message": "No files found for the specified prefix"})
            }

        files = response['Contents']
        collection_name = f"coll_{user_id}_{context_id}" if len(files) > 1 else f"one_{user_id}_{context_id}"
        vectorstore = create_vectorstore(collection_name)

        combined_embeddings = []
        combined_chunks = []
        
        # Process each file
        for file in files:
            s3_key = file['Key']
            text = read_file_from_s3(bucket_name, s3_key)
            
            if text:
                embeddings, chunks = generate_embeddings(text)
                combined_embeddings.extend(embeddings)
                combined_chunks.extend(chunks)

        # Store all embeddings at once
        metadata = {
            "context_id": context_id,
            "user_id": user_id,
            "file_count": len(files)
        }
        store_embeddings_in_qdrant(vectorstore, combined_embeddings, combined_chunks, metadata)
        
        return {
            "statusCode": 200,
            "body": json.dumps({"message": "Embeddings processed and stored successfully"})
        }

    except Exception as e:
        logger.error(f"Error in lambda_handler: {e}")
        return {
            "statusCode": 500,
            "body": json.dumps({"message": "An error occurred"})
        }
